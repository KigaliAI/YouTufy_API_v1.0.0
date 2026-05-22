from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Colors from your template
BLACK = RGBColor(0, 0, 0)
YELLOW = RGBColor(255, 215, 0)
BLUE = RGBColor(0, 112, 192)
WHITE = RGBColor(255, 255, 255)

def add_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

    # Title bar
    title_box = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.2), Inches(12), Inches(1)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = BLUE
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.color.rgb = WHITE
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Content box
    content_box = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)
    )
    content_box.line.color.rgb = BLUE
    content_box.line.width = Pt(3)
    content_tf = content_box.text_frame

    for bullet in bullets:
        p = content_tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(28)
        p.font.color.rgb = YELLOW
        p.level = 0

# Slides
add_slide("YouTufy", [
    "Your Personalized YouTube Dashboard",
    "Google OAuth 2.0 • FastAPI • Netlify • Docker"
])

add_slide("What is YouTufy?", [
    "A personalized dashboard for your YouTube subscriptions",
    "Secure Google OAuth login",
    "Clean analytics and channel insights"
])

add_slide("Key Features", [
    "Google OAuth 2.0 authentication",
    "Fetch YouTube subscriptions",
    "Dashboard with analytics",
    "Dockerized backend",
    "JWT-secured API"
])

add_slide("System Architecture", [
    "Frontend: Netlify static SPA",
    "Backend: FastAPI on Ubuntu VM",
    "Database: PostgreSQL + JSON storage",
    "NGINX reverse proxy + SSL"
])

add_slide("Deployment Diagram", [
    "User → Netlify → FastAPI → Google OAuth → YouTube API"
])

add_slide("Google OAuth Flow", [
    "User clicks 'Continue with Google'",
    "Redirect to Google Consent Screen",
    "Backend exchanges auth code for tokens",
    "JWT generated and returned to frontend",
    "Subscriptions fetched securely"
])

add_slide("Backend Stack", [
    "FastAPI (Python 3.11)",
    "SQLAlchemy ORM",
    "OAuthlib + google-auth",
    "PostgreSQL",
    "Docker + NGINX + SSL"
])

add_slide("Frontend Stack", [
    "Netlify hosting",
    "HTML + CSS + JS",
    "OAuth redirect handling",
    "Responsive dashboard UI"
])

add_slide("Project Structure", [
    "api/routes – FastAPI endpoints",
    "backend – OAuth, DB, YouTube logic",
    "templates – HTML templates",
    "static – CSS, favicon, logo",
    "users – JSON storage per user"
])

add_slide("Security Model", [
    "HTTPS enforced",
    "JWT authentication",
    "Refresh token handling",
    "YouTube scope: youtube.readonly"
])

add_slide("Subscriptions Dashboard", [
    "Channel title, thumbnail, stats",
    "Latest upload date",
    "Clean overview of all subscriptions"
])

add_slide("Data Pipeline", [
    "User logs in via Google",
    "Backend retrieves subscriptions",
    "Data stored in JSON per user",
    "Dashboard fetches via JWT"
])

add_slide("Deployment Summary", [
    "Frontend: https://www.youtufy.com",
    "Backend: https://api.youtufy.com",
    "DNS: youtufy.com → www.youtufy.com"
])

add_slide("Future Improvements", [
    "AI-powered insights",
    "Mobile app version",
    "Channel growth tracking",
    "Watch-history analytics"
])

add_slide("Thank You", [
    "Questions?",
    "YouTufy – Your Personalized YouTube Dashboard"
])

prs.save("YouTufy_Presentation.pptx")
